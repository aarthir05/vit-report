"""
generate_ppt_templates.py  v2
6 cute PowerPoint templates — 10 slides each.
All positions use Inches(); no coordinate overflow.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os, random

ROOT = r"C:\Users\VJ872NS\OneDrive - EY\Desktop\VIT_Report"
OUT  = os.path.join(ROOT, "products", "templates", "dist")
os.makedirs(OUT, exist_ok=True)

W, H = Inches(13.33), Inches(7.5)

def rgb(r, g, b): return RGBColor(r, g, b)

def rect(sl, l, t, w, h, fill):
    s = sl.shapes.add_shape(1, l, t, w, h)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def rrect(sl, l, t, w, h, fill, border=None):
    s = sl.shapes.add_shape(5, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def circle(sl, cx, cy, r, fill):
    s = sl.shapes.add_shape(9, cx - r, cy - r, r * 2, r * 2)
    s.line.fill.background()
    s.fill.solid(); s.fill.fore_color.rgb = fill
    return s

def txt(sl, text, l, t, w, h, size, bold=False, italic=False,
        color=rgb(255, 255, 255), align=PP_ALIGN.LEFT, font="Calibri"):
    tx = sl.shapes.add_textbox(l, t, w, h)
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    rn = p.add_run()
    rn.text = text; rn.font.size = Pt(size)
    rn.font.bold = bold; rn.font.italic = italic
    rn.font.color.rgb = color; rn.font.name = font
    return tx

def hline(sl, y, color, x=Inches(0.5), w=Inches(12.33)):
    rect(sl, x, y, w, Inches(0.045), color)

def blank(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def new_prs():
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H; return prs

def _three_cards(sl, cards, card_fill, hdr_color, title_color, body_color, font_h, font_b, border=None):
    for i, (title, body) in enumerate(cards):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), card_fill, border)
        txt(sl, title, x + Inches(0.2), Inches(1.85), Inches(3.55), Inches(0.75), 20, bold=True, color=title_color, font=font_h)
        hline(sl, Inches(2.68), hdr_color, x + Inches(0.2), Inches(3.55))
        txt(sl, body, x + Inches(0.2), Inches(2.8), Inches(3.55), Inches(4.15), 16, color=body_color, font=font_b)

def _stats_row(sl, stats, card_fill, num_color, lbl_color, sub_color, font_h, font_b, border=None):
    for i, (num, label, sub) in enumerate(stats):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), card_fill, border)
        txt(sl, num,   x + Inches(0.2), Inches(2.3),  Inches(3.55), Inches(1.65), 64, bold=True, color=num_color, font=font_h, align=PP_ALIGN.CENTER)
        txt(sl, label, x + Inches(0.2), Inches(4.05), Inches(3.55), Inches(0.75), 22, bold=True, color=lbl_color, font=font_h, align=PP_ALIGN.CENTER)
        txt(sl, sub,   x + Inches(0.2), Inches(4.88), Inches(3.55), Inches(0.65), 14, color=sub_color, font=font_b, align=PP_ALIGN.CENTER)

def _two_col(sl, left_fill, right_fill, left_accent, right_accent, title_color, body_color, font_h, font_b):
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(5.65), left_fill)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(5.65), right_fill)
    txt(sl, "Left Column",  Inches(0.75), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=title_color, font=font_h)
    txt(sl, "Right Column", Inches(7.08), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=title_color, font=font_h)
    hline(sl, Inches(2.52), left_accent,  Inches(0.75), Inches(5.5))
    hline(sl, Inches(2.52), right_accent, Inches(7.08), Inches(5.5))
    items = ["- Point one", "- Point two", "- Point three", "- Point four"]
    for i, item in enumerate(items):
        txt(sl, item, Inches(0.75), Inches(2.72 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=body_color, font=font_b)
        txt(sl, item, Inches(7.08), Inches(2.72 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=body_color, font=font_b)

CARDS_TEMPLATE = [
    ("Topic One",   "Add your description here. Summarise the key idea clearly and concisely."),
    ("Topic Two",   "Add your description here. Summarise the key idea clearly and concisely."),
    ("Topic Three", "Add your description here. Summarise the key idea clearly and concisely."),
]
STATS_TEMPLATE = [
    ("97%",  "Accuracy",  "Across all test cases"),
    ("24/7", "Available", "Round-the-clock support"),
    ("150+", "Clients",   "Globally served"),
]
BULLETS_TEMPLATE = [
    "First key insight - add a concise supporting statement here.",
    "Second key insight - add a concise supporting statement here.",
    "Third key insight - add a concise supporting statement here.",
    "Fourth key insight - add a concise supporting statement here.",
]
INTRO_BODY = ("About This Presentation\n\nAdd a brief introduction here. "
              "Describe the purpose, scope, and key themes of your work.\n\n"
              "Give your audience the context they need before the main content.")

# ====== THEME 1: PASTEL STUDY ======
def theme_pastel():
    CREAM = rgb(255, 253, 240); PINK  = rgb(255, 182, 193)
    LAVEN = rgb(216, 191, 216); MINT  = rgb(180, 235, 210)
    DPINK = rgb(200,  90, 120); DARK  = rgb( 70,  50,  60)

    def deco(sl):
        circle(sl, Inches(12.8),  Inches(0.9),  Inches(1.8), PINK)
        circle(sl, Inches(11.55), Inches(6.9),  Inches(1.3), MINT)
        circle(sl, Inches(0.45),  Inches(6.65), Inches(1.0), LAVEN)
        circle(sl, Inches(0.95),  Inches(1.25), Inches(0.65), PINK)

    def head(sl, title, bg=LAVEN):
        rect(sl, 0, 0, W, Inches(1.4), bg)
        txt(sl, title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=DARK, font="Georgia")

    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); deco(sl)
    rect(sl, Inches(1.2), Inches(2.6), Inches(0.1), Inches(2.35), DPINK)
    txt(sl, "Your Presentation Title", Inches(1.5), Inches(2.6), Inches(9.5), Inches(1.45), 46, bold=True, color=DARK, font="Georgia")
    txt(sl, "Subtitle  -  Presenter Name  -  Date", Inches(1.5), Inches(4.2), Inches(9.5), Inches(0.7), 17, color=DPINK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM)
    rect(sl, 0, 0, Inches(4.6), H, PINK)
    circle(sl, Inches(4.6), Inches(7.35), Inches(1.6), LAVEN)
    txt(sl, "Agenda", Inches(0.3), Inches(2.95), Inches(3.9), Inches(1.1), 42, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    for i, item in enumerate(["01  Introduction", "02  Key Points", "03  Discussion", "04  Next Steps", "05  Wrap Up"]):
        txt(sl, item, Inches(5.2), Inches(1.35) + Inches(i * 0.95), Inches(7.5), Inches(0.85), 20, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Introduction", MINT)
    rrect(sl, Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.55), PINK)
    txt(sl, "[ Image / Visual Placeholder ]", Inches(0.5), Inches(3.9), Inches(4.5), Inches(0.7), 13, color=DARK, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, INTRO_BODY, Inches(5.35), Inches(1.65), Inches(7.5), Inches(5.35), 17, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, DPINK)
    circle(sl, Inches(0.8),  Inches(0.9),  Inches(2.2), rgb(215, 100, 130))
    circle(sl, Inches(12.6), Inches(6.75), Inches(2.4), rgb(215, 100, 130))
    circle(sl, Inches(13.1), Inches(0.6),  Inches(1.2), rgb(215, 100, 130))
    txt(sl, "Chapter 01", Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.8), 22, color=rgb(255, 220, 230), font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, "Section Title Here", Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.35), 52, bold=True, color=rgb(255, 245, 248), font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.75), rgb(255, 220, 230), Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Key Topics", PINK)
    _three_cards(sl, CARDS_TEMPLATE, PINK, DPINK, DARK, DARK, "Georgia", "Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Key Takeaways", LAVEN)
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        circle(sl, Inches(0.95), y + Inches(0.3), Inches(0.22), DPINK)
        txt(sl, item, Inches(1.35), y, Inches(11.0), Inches(1.1), 19, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, PINK)
    circle(sl, Inches(0.9), Inches(1.0), Inches(1.8), rgb(255, 200, 210))
    circle(sl, Inches(12.45), Inches(6.5), Inches(2.0), LAVEN)
    txt(sl, "\u201c", Inches(0.5), Inches(1.1), Inches(2.5), Inches(2.0), 120, bold=True, color=DPINK, font="Georgia")
    txt(sl, "Add your most impactful quote or key message here. Make it memorable.", Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.1), 34, italic=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "- Attribution Name, Role", Inches(1.2), Inches(4.65), Inches(10.8), Inches(0.65), 17, color=DPINK, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "By the Numbers", MINT)
    _stats_row(sl, STATS_TEMPLATE, LAVEN, DPINK, DARK, DARK, "Georgia", "Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Two Column Layout", PINK)
    _two_col(sl, PINK, LAVEN, DPINK, DPINK, DARK, DARK, "Georgia", "Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, PINK); deco(sl)
    txt(sl, "Thank You!", Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.45), 52, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.15), DPINK, Inches(3.0), Inches(7.33))
    txt(sl, "yourname@email.com  -  @handle  -  website.com", Inches(1.0), Inches(4.35), Inches(11.3), Inches(0.7), 17, color=DPINK, font="Calibri", align=PP_ALIGN.CENTER)

    prs.save(os.path.join(OUT, "01_pastel_study.pptx")); print("[OK] 01_pastel_study.pptx")

# ====== THEME 2: NOTION MINIMAL ======
def theme_notion():
    CREAM = rgb(250, 247, 240); DARK  = rgb( 37,  34,  28)
    OLIVE = rgb(100, 115,  70); TERRA = rgb(185,  95,  60)
    STONE = rgb(210, 200, 185); LTOLV = rgb(220, 230, 200)
    TAN   = rgb(238, 228, 210)

    def dotgrid(sl):
        for gx in range(23):
            for gy in range(13):
                circle(sl, Inches(gx * 0.6), Inches(gy * 0.6), Inches(0.03), rgb(215, 208, 195))

    def head(sl, title):
        rect(sl, 0, 0, W, Inches(1.4), LTOLV)
        rect(sl, Inches(0.5), Inches(0.28), Inches(0.07), Inches(0.85), TERRA)
        txt(sl, title, Inches(0.72), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=DARK, font="Georgia")

    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl)
    rect(sl, Inches(11.0), 0, Inches(2.33), H, LTOLV)
    rect(sl, Inches(0.8), Inches(3.25), Inches(4.2), Inches(0.055), TERRA)
    txt(sl, "Your Presentation", Inches(0.8), Inches(3.4), Inches(9.5), Inches(1.35), 46, bold=True, color=DARK, font="Georgia")
    txt(sl, "A minimal template  -  Author  -  Date", Inches(0.8), Inches(4.85), Inches(9.5), Inches(0.7), 17, color=OLIVE, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl)
    rect(sl, Inches(0.5), Inches(0.8), Inches(0.065), Inches(1.0), TERRA)
    txt(sl, "Agenda", Inches(0.75), Inches(0.7), Inches(11.5), Inches(0.95), 36, bold=True, color=DARK, font="Georgia")
    hline(sl, Inches(1.85), STONE)
    for i, item in enumerate(["Introduction","Background & Context","Main Findings","Discussion","Conclusion"]):
        y = Inches(2.15 + i * 0.88)
        rect(sl, Inches(0.8), y + Inches(0.2), Inches(0.26), Inches(0.26), TERRA)
        txt(sl, "  " + item, Inches(1.25), y, Inches(11.0), Inches(0.8), 19, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl)
    rect(sl, 0, 0, Inches(3.8), H, LTOLV)
    rect(sl, Inches(0.5), Inches(0.7), Inches(0.065), Inches(0.95), TERRA)
    txt(sl, "Introduction", Inches(0.75), Inches(0.65), Inches(3.0), Inches(0.95), 26, bold=True, color=DARK, font="Georgia")
    txt(sl, INTRO_BODY, Inches(4.2), Inches(1.3), Inches(8.7), Inches(5.8), 17, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, LTOLV)
    txt(sl, "Chapter 01", Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.8), 20, color=TERRA, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, "Section Title Here", Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.35), 52, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.75), TERRA, Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl); head(sl, "Key Topics")
    _three_cards(sl, CARDS_TEMPLATE, TAN, TERRA, DARK, DARK, "Georgia", "Calibri", STONE)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl); head(sl, "Key Takeaways")
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        txt(sl, "0" + str(i+1), Inches(0.7), y, Inches(0.5), Inches(0.55), 22, bold=True, color=TERRA, font="Georgia")
        txt(sl, item, Inches(1.35), y, Inches(11.0), Inches(1.1), 18, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl)
    rect(sl, Inches(0.5), Inches(1.5), Inches(0.12), Inches(4.5), TERRA)
    txt(sl, "\u201cAdd your most impactful quote here. Make it memorable.\u201d", Inches(1.0), Inches(2.15), Inches(11.5), Inches(2.5), 32, italic=True, color=DARK, font="Georgia")
    txt(sl, "- Attribution Name, Role", Inches(1.0), Inches(4.85), Inches(11.5), Inches(0.65), 17, color=OLIVE, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl); head(sl, "By the Numbers")
    _stats_row(sl, STATS_TEMPLATE, TAN, TERRA, DARK, DARK, "Georgia", "Calibri", STONE)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl); head(sl, "Two Column Layout")
    _two_col(sl, LTOLV, TAN, TERRA, OLIVE, DARK, DARK, "Georgia", "Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); dotgrid(sl)
    rect(sl, 0, 0, Inches(3.8), H, LTOLV)
    rect(sl, Inches(3.8), Inches(3.6), Inches(9.53), Inches(0.055), TERRA)
    txt(sl, "Thank\nYou", Inches(0.3), Inches(2.5), Inches(3.1), Inches(2.0), 44, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "yourname@email.com\n@handle  -  website.com", Inches(4.3), Inches(3.8), Inches(8.7), Inches(1.1), 17, color=OLIVE, font="Calibri")

    prs.save(os.path.join(OUT, "02_notion_minimal.pptx")); print("[OK] 02_notion_minimal.pptx")

# ====== THEME 3: NIGHT SKY ======
def theme_nightsky():
    NAVY  = rgb( 10,  15,  45); DNAVY = rgb(  5,   8,  30)
    GOLD  = rgb(255, 215,  80); LGOLD = rgb(255, 235, 140)
    BLUE  = rgb( 80, 120, 220); LBLUE = rgb(140, 180, 255)
    WHITE = rgb(255, 255, 255); WDIM  = rgb(170, 190, 230)
    PANEL = rgb( 15,  22,  65)

    def stars(sl, seed=42, count=65):
        random.seed(seed)
        for _ in range(count):
            sx = random.uniform(0.1, 13.2); sy = random.uniform(0.1, 7.4)
            sr = random.uniform(0.025, 0.065); br = random.randint(170, 255)
            circle(sl, Inches(sx), Inches(sy), Inches(sr), rgb(br, br, min(255, br + 40)))

    def cloud(sl, cx, cy, sc, col):
        for ox, oy, r in [(-sc*0.5, 0, sc*0.35), (0, -sc*0.2, sc*0.5), (sc*0.5, 0, sc*0.35), (0, sc*0.15, sc*0.4)]:
            circle(sl, Inches(cx + ox), Inches(cy + oy), Inches(r), col)

    def head(sl, title):
        rect(sl, 0, 0, W, Inches(1.4), PANEL)
        txt(sl, "* " + title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=GOLD, font="Georgia")

    ACCS = [GOLD, BLUE, LGOLD]
    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 1)
    circle(sl, Inches(11.5), Inches(1.2), Inches(1.6), rgb(25, 35, 80))
    circle(sl, Inches(11.85), Inches(0.9), Inches(1.6), rgb(60, 75, 160))
    circle(sl, Inches(11.95), Inches(0.85), Inches(1.4), NAVY)
    cloud(sl, 1.3, 6.55, 1.3, rgb(20, 28, 70)); cloud(sl, 11.8, 6.3, 1.1, rgb(20, 28, 70))
    rect(sl, Inches(0.8), Inches(3.75), Inches(4.5), Inches(0.055), GOLD)
    txt(sl, "Your Presentation", Inches(0.8), Inches(2.35), Inches(10.5), Inches(1.35), 46, bold=True, color=WHITE, font="Georgia")
    txt(sl, "Subtitle  -  Name  -  Date", Inches(0.8), Inches(3.85), Inches(8.5), Inches(0.7), 18, color=LBLUE, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 2)
    rect(sl, 0, 0, Inches(4.6), H, PANEL); cloud(sl, 6.7, 6.55, 1.1, rgb(20, 28, 70))
    txt(sl, "* Agenda", Inches(0.3), Inches(2.95), Inches(3.9), Inches(1.1), 34, bold=True, color=GOLD, font="Georgia", align=PP_ALIGN.CENTER)
    for i, item in enumerate(["01  Introduction","02  Key Points","03  Discussion","04  Next Steps","05  Wrap Up"]):
        txt(sl, item, Inches(5.2), Inches(1.35) + Inches(i * 0.95), Inches(7.5), Inches(0.85), 20, color=WDIM, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 3)
    rect(sl, Inches(8.33), 0, Inches(5.0), H, PANEL); head(sl, "Introduction")
    txt(sl, INTRO_BODY, Inches(0.55), Inches(1.65), Inches(7.45), Inches(5.4), 17, color=WDIM, font="Calibri")
    txt(sl, "[ Visual / Image ]", Inches(8.65), Inches(3.5), Inches(4.35), Inches(0.7), 13, color=LBLUE, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, DNAVY); stars(sl, 4, count=90)
    txt(sl, "Chapter 01", Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.8), 20, color=LBLUE, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, "Section Title Here", Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.35), 52, bold=True, color=GOLD, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.75), GOLD, Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 5); head(sl, "Key Topics")
    for i, (title, body) in enumerate(CARDS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), PANEL)
        rect(sl, x, Inches(1.6), Inches(3.97), Inches(0.07), ACCS[i])
        txt(sl, title, x + Inches(0.2), Inches(1.85), Inches(3.55), Inches(0.75), 20, bold=True, color=GOLD, font="Georgia")
        txt(sl, body, x + Inches(0.2), Inches(2.75), Inches(3.55), Inches(4.2), 16, color=WDIM, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 6); head(sl, "Key Takeaways")
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        txt(sl, "*", Inches(0.6), y, Inches(0.6), Inches(0.6), 22, bold=True, color=GOLD, font="Georgia")
        txt(sl, item, Inches(1.35), y, Inches(11.0), Inches(1.1), 18, color=WDIM, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 7)
    txt(sl, "\u201c", Inches(0.5), Inches(1.1), Inches(2.5), Inches(2.0), 120, bold=True, color=GOLD, font="Georgia")
    txt(sl, "Add your most impactful quote or key message here. Make it memorable.", Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.1), 34, italic=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "- Attribution Name, Role", Inches(1.2), Inches(4.65), Inches(10.8), Inches(0.65), 17, color=LBLUE, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 8); head(sl, "By the Numbers")
    for i, (num, label, sub) in enumerate(STATS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), PANEL)
        rect(sl, x, Inches(1.6), Inches(3.97), Inches(0.07), ACCS[i])
        txt(sl, num,   x + Inches(0.2), Inches(2.3),  Inches(3.55), Inches(1.65), 64, bold=True, color=GOLD, font="Georgia", align=PP_ALIGN.CENTER)
        txt(sl, label, x + Inches(0.2), Inches(4.05), Inches(3.55), Inches(0.75), 22, bold=True, color=WHITE, font="Georgia", align=PP_ALIGN.CENTER)
        txt(sl, sub,   x + Inches(0.2), Inches(4.88), Inches(3.55), Inches(0.65), 14, color=WDIM, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, NAVY); stars(sl, 9); head(sl, "Two Column Layout")
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(5.65), PANEL)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(5.65), PANEL)
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(0.07), BLUE)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(0.07), GOLD)
    txt(sl, "Left Column",  Inches(0.75), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=BLUE, font="Georgia")
    txt(sl, "Right Column", Inches(7.08), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=GOLD, font="Georgia")
    for i, item in enumerate(["- Point one","- Point two","- Point three","- Point four"]):
        txt(sl, item, Inches(0.75), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=WDIM, font="Calibri")
        txt(sl, item, Inches(7.08), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=WDIM, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, DNAVY); stars(sl, 10, count=80)
    circle(sl, Inches(6.67), Inches(2.3), Inches(2.6), rgb(15, 25, 70))
    cloud(sl, 1.8, 6.35, 1.4, rgb(15, 22, 60)); cloud(sl, 11.2, 6.25, 1.2, rgb(15, 22, 60))
    txt(sl, "* Thank You *", Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.35), 52, bold=True, color=GOLD, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.1), GOLD, Inches(3.0), Inches(7.33))
    txt(sl, "yourname@email.com  -  @handle  -  website.com", Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.7), 17, color=LBLUE, font="Calibri", align=PP_ALIGN.CENTER)

    prs.save(os.path.join(OUT, "03_nightsky.pptx")); print("[OK] 03_nightsky.pptx")

# ====== THEME 4: COTTAGECORE BOHO ======
def theme_cottagecore():
    BEIGE  = rgb(245, 235, 215); DBROWN = rgb( 90,  58,  30)
    TERRA  = rgb(195, 105,  65); GREEN  = rgb( 90, 125,  75)
    SAGE   = rgb(185, 205, 175); WARM   = rgb(230, 205, 170)

    def leaf_cluster(sl, cx, cy, sc=1.0):
        for ox, oy, ew, eh, col in [(0, 0, sc*0.65, sc*1.5, GREEN), (sc*0.4, -sc*0.2, sc*0.45, sc*1.1, SAGE), (-sc*0.35, sc*0.1, sc*0.35, sc*0.9, GREEN)]:
            s = sl.shapes.add_shape(9, Inches(cx+ox-ew/2), Inches(cy+oy-eh/2), Inches(ew), Inches(eh))
            s.line.fill.background(); s.fill.solid(); s.fill.fore_color.rgb = col

    def head(sl, title, bg=WARM):
        rect(sl, 0, 0, W, Inches(1.4), bg)
        txt(sl, title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=DBROWN, font="Georgia")

    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE)
    for i in range(8):
        r = Inches(0.45 + i * 0.25)
        s = sl.shapes.add_shape(9, -r, -r, r * 2, r * 2)
        s.fill.background(); s.line.color.rgb = SAGE; s.line.width = Pt(1)
    for i in range(6):
        r = Inches(0.4 + i * 0.22)
        s = sl.shapes.add_shape(9, W - r, H - r, r * 2, r * 2)
        s.fill.background(); s.line.color.rgb = TERRA; s.line.width = Pt(1)
    leaf_cluster(sl, 12.7, 0.7, 1.0); leaf_cluster(sl, 0.6, 6.9, 0.8)
    rect(sl, Inches(1.5), Inches(3.65), Inches(3.2), Inches(0.055), TERRA)
    txt(sl, "Your Presentation Title", Inches(1.5), Inches(2.4), Inches(9.5), Inches(1.3), 46, bold=True, color=DBROWN, font="Georgia")
    txt(sl, "Subtitle  -  Author Name  -  Date", Inches(1.5), Inches(3.8), Inches(9.5), Inches(0.7), 17, color=TERRA, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); rect(sl, 0, 0, Inches(3.8), H, WARM)
    leaf_cluster(sl, 3.4, 0.5, 0.8)
    txt(sl, "Agenda", Inches(0.3), Inches(3.0), Inches(3.1), Inches(1.0), 36, bold=True, color=DBROWN, font="Georgia", align=PP_ALIGN.CENTER)
    for i, item in enumerate(["Introduction","Background","Main Content","Discussion","Closing Notes"]):
        rect(sl, Inches(4.4), Inches(1.35 + i * 0.95) + Inches(0.22), Inches(0.22), Inches(0.22), TERRA)
        txt(sl, item, Inches(4.85), Inches(1.35 + i * 0.95), Inches(8.0), Inches(0.85), 20, color=DBROWN, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); head(sl, "Introduction")
    rrect(sl, Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.55), WARM, TERRA)
    txt(sl, "[ Image / Visual ]", Inches(0.5), Inches(3.9), Inches(4.5), Inches(0.7), 13, color=DBROWN, font="Calibri", align=PP_ALIGN.CENTER)
    leaf_cluster(sl, 4.85, 1.6, 0.6)
    txt(sl, INTRO_BODY, Inches(5.35), Inches(1.65), Inches(7.5), Inches(5.35), 17, color=DBROWN, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, WARM)
    leaf_cluster(sl, 0.9, 0.9, 1.1); leaf_cluster(sl, 12.5, 6.7, 1.0)
    txt(sl, "Chapter 01", Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.8), 20, color=TERRA, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, "Section Title Here", Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.35), 52, bold=True, color=DBROWN, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.75), TERRA, Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); head(sl, "Key Topics")
    _three_cards(sl, CARDS_TEMPLATE, WARM, TERRA, DBROWN, DBROWN, "Georgia", "Calibri", TERRA)

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); head(sl, "Key Takeaways")
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        rect(sl, Inches(0.72), y + Inches(0.22), Inches(0.24), Inches(0.24), TERRA)
        txt(sl, item, Inches(1.2), y, Inches(11.3), Inches(1.1), 19, color=DBROWN, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, SAGE)
    leaf_cluster(sl, 0.9, 0.9, 0.9); leaf_cluster(sl, 12.5, 6.6, 0.8)
    txt(sl, "\u201c", Inches(0.5), Inches(1.1), Inches(2.5), Inches(2.0), 120, bold=True, color=TERRA, font="Georgia")
    txt(sl, "Add your most impactful quote or key message here. Make it memorable.", Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.1), 34, italic=True, color=DBROWN, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "- Attribution Name, Role", Inches(1.2), Inches(4.65), Inches(10.8), Inches(0.65), 17, color=TERRA, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); head(sl, "By the Numbers")
    _stats_row(sl, STATS_TEMPLATE, WARM, TERRA, DBROWN, DBROWN, "Georgia", "Calibri", TERRA)

    sl = blank(prs); rect(sl, 0, 0, W, H, BEIGE); head(sl, "Two Column Layout")
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(5.65), WARM)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(5.65), rgb(235, 225, 205))
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(0.1), TERRA)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(0.1), GREEN)
    txt(sl, "Left Column",  Inches(0.75), Inches(1.8), Inches(5.5), Inches(0.65), 22, bold=True, color=DBROWN, font="Georgia")
    txt(sl, "Right Column", Inches(7.08), Inches(1.8), Inches(5.5), Inches(0.65), 22, bold=True, color=DBROWN, font="Georgia")
    for i, item in enumerate(["- Point one","- Point two","- Point three","- Point four"]):
        txt(sl, item, Inches(0.75), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=DBROWN, font="Calibri")
        txt(sl, item, Inches(7.08), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=DBROWN, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, WARM)
    for i in range(10):
        r = Inches(0.45 + i * 0.28)
        s = sl.shapes.add_shape(9, W - r, H - r, r * 2, r * 2)
        s.fill.background(); s.line.color.rgb = SAGE; s.line.width = Pt(1.2)
    leaf_cluster(sl, 1.1, 1.0, 1.0); leaf_cluster(sl, 12.4, 6.7, 0.9)
    txt(sl, "Thank You", Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.35), 52, bold=True, color=DBROWN, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.1), TERRA, Inches(3.0), Inches(7.33))
    txt(sl, "yourname@email.com  -  @handle  -  website.com", Inches(1.0), Inches(4.3), Inches(11.3), Inches(0.7), 17, color=TERRA, font="Calibri", align=PP_ALIGN.CENTER)

    prs.save(os.path.join(OUT, "04_cottagecore_boho.pptx")); print("[OK] 04_cottagecore_boho.pptx")

# ====== THEME 5: Y2K RETRO ======
def theme_y2k():
    BG   = rgb( 12,   5,  28); PINK = rgb(255,  60, 150)
    CYAN = rgb(  0, 220, 220); YELL = rgb(255, 230,  30)
    PURP = rgb(160,  40, 220); WHITE= rgb(255, 255, 255)

    def head(sl, title, bg=PURP):
        rect(sl, 0, 0, W, Inches(1.4), bg)
        txt(sl, title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=WHITE, font="Arial Black")

    ACCS = [PINK, CYAN, YELL]
    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, BG)
    circle(sl, Inches(1.5), Inches(1.2), Inches(1.2), PURP)
    circle(sl, Inches(11.8), Inches(6.25), Inches(1.5), CYAN)
    circle(sl, Inches(12.5), Inches(1.5), Inches(0.8), PINK)
    circle(sl, Inches(0.8), Inches(6.5), Inches(0.6), YELL)
    rect(sl, Inches(0.7), Inches(2.15), Inches(10.0), Inches(1.5), PINK)
    txt(sl, "YOUR TITLE HERE", Inches(0.8), Inches(2.2), Inches(9.8), Inches(1.35), 52, bold=True, color=BG, font="Arial Black", align=PP_ALIGN.CENTER)
    txt(sl, "Subtitle  *  Name  *  Date", Inches(0.7), Inches(3.9), Inches(10.0), Inches(0.75), 20, color=CYAN, font="Arial", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, BG)
    circle(sl, Inches(12.8), Inches(0.65), Inches(1.8), PURP)
    circle(sl, Inches(0.5), Inches(7.1), Inches(1.2), CYAN)
    rect(sl, Inches(0.6), Inches(0.5), Inches(3.5), Inches(1.0), PINK)
    txt(sl, "AGENDA", Inches(0.6), Inches(0.55), Inches(3.5), Inches(0.9), 34, bold=True, color=BG, font="Arial Black", align=PP_ALIGN.CENTER)
    item_cols = [CYAN, PINK, YELL, PURP, CYAN]
    for i, item in enumerate(["* Introduction","* Key Points","* Discussion","* Next Steps","* Wrap Up"]):
        txt(sl, item, Inches(1.2), Inches(1.8 + i * 0.96), Inches(10.0), Inches(0.85), 22, color=item_cols[i], bold=True, font="Arial")

    sl = blank(prs); rect(sl, 0, 0, W, H, BG)
    circle(sl, Inches(12.2), Inches(1.5), Inches(1.4), PURP)
    circle(sl, Inches(1.0), Inches(6.5), Inches(0.9), CYAN)
    head(sl, "INTRODUCTION", PURP)
    rrect(sl, Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.55), rgb(30, 15, 60))
    txt(sl, "[ Visual / Image ]", Inches(0.5), Inches(3.9), Inches(4.5), Inches(0.7), 13, color=CYAN, font="Arial", align=PP_ALIGN.CENTER)
    txt(sl, INTRO_BODY, Inches(5.35), Inches(1.65), Inches(7.5), Inches(5.35), 17, color=WHITE, font="Arial")

    sl = blank(prs); rect(sl, 0, 0, W, H, PINK)
    circle(sl, Inches(0.8), Inches(0.9), Inches(1.8), rgb(230, 40, 130))
    circle(sl, Inches(12.6), Inches(6.75), Inches(2.0), rgb(230, 40, 130))
    txt(sl, "CHAPTER 01", Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.8), 22, color=BG, font="Arial Black", align=PP_ALIGN.CENTER)
    txt(sl, "SECTION TITLE HERE", Inches(0.5), Inches(3.25), Inches(12.3), Inches(1.35), 52, bold=True, color=BG, font="Arial Black", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.75), BG, Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, BG); head(sl, "KEY TOPICS", PURP)
    for i, (title, body) in enumerate(CARDS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), ACCS[i])
        txt(sl, "CARD 0" + str(i+1), x + Inches(0.2), Inches(1.85), Inches(3.55), Inches(0.75), 22, bold=True, color=BG, font="Arial Black")
        txt(sl, body, x + Inches(0.2), Inches(2.75), Inches(3.55), Inches(4.2), 16, color=BG, font="Arial")

    sl = blank(prs); rect(sl, 0, 0, W, H, BG); head(sl, "KEY TAKEAWAYS", CYAN)
    bul_cols = [PINK, CYAN, YELL, PURP]
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        txt(sl, "*", Inches(0.6), y, Inches(0.6), Inches(0.6), 24, bold=True, color=bul_cols[i], font="Arial")
        txt(sl, item, Inches(1.35), y, Inches(11.0), Inches(1.1), 18, color=WHITE, font="Arial")

    sl = blank(prs); rect(sl, 0, 0, W, H, BG)
    circle(sl, Inches(0.9), Inches(1.0), Inches(1.5), PURP)
    circle(sl, Inches(12.4), Inches(6.5), Inches(1.8), CYAN)
    txt(sl, "\u201c", Inches(0.5), Inches(1.1), Inches(2.5), Inches(2.0), 120, bold=True, color=PINK, font="Arial Black")
    txt(sl, "Add your most impactful quote or key message here. Make it memorable.", Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.1), 34, italic=True, color=WHITE, font="Arial", align=PP_ALIGN.CENTER)
    txt(sl, "- Attribution Name, Role", Inches(1.2), Inches(4.65), Inches(10.8), Inches(0.65), 17, color=YELL, font="Arial", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, BG); head(sl, "BY THE NUMBERS", PURP)
    for i, (num, label, sub) in enumerate(STATS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), rgb(30, 15, 60))
        rect(sl, x, Inches(1.6), Inches(3.97), Inches(0.08), ACCS[i])
        txt(sl, num,   x + Inches(0.2), Inches(2.3),  Inches(3.55), Inches(1.65), 64, bold=True, color=ACCS[i], font="Arial Black", align=PP_ALIGN.CENTER)
        txt(sl, label.upper(), x + Inches(0.2), Inches(4.05), Inches(3.55), Inches(0.75), 22, bold=True, color=WHITE, font="Arial Black", align=PP_ALIGN.CENTER)
        txt(sl, sub, x + Inches(0.2), Inches(4.88), Inches(3.55), Inches(0.65), 13, color=rgb(200, 190, 220), font="Arial", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, BG); head(sl, "TWO COLUMN LAYOUT", CYAN)
    rect(sl, Inches(0.5),  Inches(1.55), Inches(6.0), Inches(5.65), PINK)
    rect(sl, Inches(6.83), Inches(1.55), Inches(6.0), Inches(5.65), PURP)
    txt(sl, "Left Column",  Inches(0.75), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=BG, font="Arial Black")
    txt(sl, "Right Column", Inches(7.08), Inches(1.75), Inches(5.5), Inches(0.65), 22, bold=True, color=WHITE, font="Arial Black")
    for i, item in enumerate(["* Point one","* Point two","* Point three","* Point four"]):
        txt(sl, item, Inches(0.75), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=BG, bold=True, font="Arial")
        txt(sl, item, Inches(7.08), Inches(2.6 + i * 0.65), Inches(5.5), Inches(0.6), 17, color=WHITE, font="Arial")

    sl = blank(prs); rect(sl, 0, 0, W, H, BG)
    circle(sl, Inches(2.0), Inches(1.5), Inches(2.2), PURP)
    circle(sl, Inches(11.5), Inches(5.8), Inches(2.0), CYAN)
    txt(sl, "*", Inches(5.5), Inches(0.8), Inches(2.33), Inches(1.2), 64, bold=True, color=YELL, font="Arial", align=PP_ALIGN.CENTER)
    rect(sl, Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.5), PINK)
    txt(sl, "THANK YOU! *", Inches(0.5), Inches(3.05), Inches(12.3), Inches(1.35), 50, bold=True, color=BG, font="Arial Black", align=PP_ALIGN.CENTER)
    txt(sl, "yourname@email.com  *  @handle  *  website.com", Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.75), 18, color=CYAN, font="Arial", align=PP_ALIGN.CENTER)

    prs.save(os.path.join(OUT, "05_y2k_retro.pptx")); print("[OK] 05_y2k_retro.pptx")

# ====== THEME 6: SANRIO BEAR ======
def theme_sanrio():
    CREAM = rgb(255, 252, 245); SOFT_P = rgb(255, 220, 230)
    YELL  = rgb(255, 240, 150); BRWN   = rgb(150,  90,  50)
    PNKDP = rgb(230, 100, 140); MID    = rgb(255, 175, 195)
    DARK  = rgb( 60,  40,  50)
    FILLS = [SOFT_P, YELL, MID]

    def polkadots(sl, seed=7, count=25):
        random.seed(seed)
        for _ in range(count):
            px = random.uniform(0.15, 13.15); py = random.uniform(0.15, 7.35)
            sz = random.uniform(0.07, 0.18)
            circle(sl, Inches(px), Inches(py), Inches(sz), random.choice(FILLS))

    def bear_ears(sl, cx, cy, sz):
        circle(sl, Inches(cx - sz*0.55), Inches(cy - sz*0.55), Inches(sz*0.40), BRWN)
        circle(sl, Inches(cx + sz*0.55), Inches(cy - sz*0.55), Inches(sz*0.40), BRWN)
        circle(sl, Inches(cx - sz*0.55), Inches(cy - sz*0.55), Inches(sz*0.25), SOFT_P)
        circle(sl, Inches(cx + sz*0.55), Inches(cy - sz*0.55), Inches(sz*0.25), SOFT_P)
        circle(sl, Inches(cx), Inches(cy), Inches(sz*0.55), BRWN)
        circle(sl, Inches(cx), Inches(cy), Inches(sz*0.42), SOFT_P)

    def head(sl, title, bg=SOFT_P):
        rect(sl, 0, 0, W, Inches(1.4), bg)
        txt(sl, title, Inches(0.55), Inches(0.28), Inches(11.5), Inches(0.85), 30, bold=True, color=DARK, font="Georgia")

    prs = new_prs()
    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); polkadots(sl, 1)
    bear_ears(sl, 12.3, 0.85, 0.95); bear_ears(sl, 1.1, 6.7, 0.75)
    rrect(sl, Inches(1.5), Inches(2.6), Inches(8.5), Inches(1.5), SOFT_P)
    txt(sl, "Your Cute Presentation!", Inches(1.55), Inches(2.65), Inches(8.4), Inches(1.35), 40, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "Subtitle  -  Name  -  Date", Inches(1.55), Inches(4.25), Inches(8.4), Inches(0.7), 18, color=PNKDP, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM)
    rect(sl, 0, 0, Inches(4.2), H, SOFT_P); polkadots(sl, 2, count=15)
    bear_ears(sl, 2.1, 0.8, 0.85)
    txt(sl, "Agenda", Inches(0.3), Inches(2.95), Inches(3.5), Inches(1.0), 34, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    for i, item in enumerate(["Introduction","Key Points","Discussion","Next Steps","Wrap Up"]):
        txt(sl, "- " + item, Inches(4.65), Inches(1.35) + Inches(i * 0.95), Inches(8.3), Inches(0.85), 20, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Introduction", YELL)
    bear_ears(sl, 12.6, 0.7, 0.8)
    rrect(sl, Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.55), SOFT_P)
    txt(sl, "[ Image / Visual ]", Inches(0.5), Inches(3.9), Inches(4.5), Inches(0.7), 13, color=DARK, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, INTRO_BODY, Inches(5.35), Inches(1.65), Inches(7.5), Inches(5.35), 17, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, SOFT_P); polkadots(sl, 4)
    bear_ears(sl, 6.67, 1.5, 1.1)
    txt(sl, "Chapter 01", Inches(0.5), Inches(2.7), Inches(12.3), Inches(0.8), 20, color=PNKDP, font="Calibri", align=PP_ALIGN.CENTER)
    txt(sl, "Section Title Here", Inches(0.5), Inches(3.4), Inches(12.3), Inches(1.3), 48, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    hline(sl, Inches(4.85), PNKDP, Inches(3.5), Inches(6.33))

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Key Topics")
    for i, (title, body) in enumerate(CARDS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), FILLS[i])
        txt(sl, title, x + Inches(0.2), Inches(1.85), Inches(3.55), Inches(0.75), 20, bold=True, color=DARK, font="Georgia")
        hline(sl, Inches(2.68), PNKDP, x + Inches(0.2), Inches(3.55))
        txt(sl, body, x + Inches(0.2), Inches(2.8), Inches(3.55), Inches(4.15), 16, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Key Takeaways")
    for i, item in enumerate(BULLETS_TEMPLATE):
        y = Inches(1.65 + i * 1.3)
        txt(sl, "-", Inches(0.65), y, Inches(0.55), Inches(0.55), 24, color=PNKDP, font="Arial", align=PP_ALIGN.CENTER)
        txt(sl, item, Inches(1.35), y, Inches(11.0), Inches(1.1), 19, color=DARK, font="Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, YELL); polkadots(sl, 7, count=15)
    txt(sl, "\u201c", Inches(0.5), Inches(1.1), Inches(2.5), Inches(2.0), 120, bold=True, color=PNKDP, font="Georgia")
    txt(sl, "Add your most impactful quote or key message here. Make it memorable.", Inches(1.2), Inches(2.4), Inches(10.8), Inches(2.1), 34, italic=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "- Attribution Name, Role", Inches(1.2), Inches(4.65), Inches(10.8), Inches(0.65), 17, color=PNKDP, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "By the Numbers")
    bear_ears(sl, 12.6, 0.7, 0.8)
    for i, (num, label, sub) in enumerate(STATS_TEMPLATE):
        x = Inches(0.5 + i * 4.27)
        rrect(sl, x, Inches(1.6), Inches(3.97), Inches(5.6), FILLS[i])
        txt(sl, num,   x + Inches(0.2), Inches(2.3),  Inches(3.55), Inches(1.65), 64, bold=True, color=PNKDP, font="Georgia", align=PP_ALIGN.CENTER)
        txt(sl, label, x + Inches(0.2), Inches(4.05), Inches(3.55), Inches(0.75), 22, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
        txt(sl, sub,   x + Inches(0.2), Inches(4.88), Inches(3.55), Inches(0.65), 14, color=DARK, font="Calibri", align=PP_ALIGN.CENTER)

    sl = blank(prs); rect(sl, 0, 0, W, H, CREAM); head(sl, "Two Column Layout")
    _two_col(sl, SOFT_P, YELL, PNKDP, PNKDP, DARK, DARK, "Georgia", "Calibri")

    sl = blank(prs); rect(sl, 0, 0, W, H, SOFT_P); polkadots(sl, 10, count=20)
    bear_ears(sl, 6.67, 1.4, 1.2)
    txt(sl, "Thank You!", Inches(0.5), Inches(3.9), Inches(12.3), Inches(1.2), 50, bold=True, color=DARK, font="Georgia", align=PP_ALIGN.CENTER)
    txt(sl, "yourname@email.com  -  @handle  -  website.com", Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.7), 17, color=PNKDP, font="Calibri", align=PP_ALIGN.CENTER)

    prs.save(os.path.join(OUT, "06_sanrio_bear.pptx")); print("[OK] 06_sanrio_bear.pptx")

if __name__ == "__main__":
    print("Generating 6 PPT templates (10 slides each)...")
    theme_pastel()
    theme_notion()
    theme_nightsky()
    theme_cottagecore()
    theme_y2k()
    theme_sanrio()
    print(f"\nAll done!  Files saved to: {OUT}")
